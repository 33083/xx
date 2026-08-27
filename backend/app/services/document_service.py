"""文档服务：上传 → 切片 → 入库（MySQL + Chroma）。"""
from __future__ import annotations

import hashlib
import os
import shutil
import uuid
from datetime import datetime
from pathlib import Path
from typing import List

from fastapi import UploadFile
from sqlalchemy import func, or_, select
from sqlalchemy.orm import Session

from app.config import settings
from app.core.vectorstore import get_user_vectorstore, user_collection_name
from app.models.document import Document
from app.schemas.document import DocCategory, DocumentOut, DocumentUploaded

_ALLOWED_EXT = {".pdf", ".txt", ".md", ".markdown", ".pptx", ".doc", ".docx"}


def _ensure_dirs() -> Path:
    p = Path(settings.DOC_STORAGE_DIR)
    p.mkdir(parents=True, exist_ok=True)
    return p


def _category_of(filename: str, hint: str | None) -> str:
    if hint and hint in {DocCategory.MATERIAL, DocCategory.RESUME, DocCategory.INTERVIEW}:
        return hint
    low = filename.lower()
    if "简历" in filename or "resume" in low or "cv" in low:
        return DocCategory.RESUME
    if "面经" in filename or "interview" in low:
        return DocCategory.INTERVIEW
    return DocCategory.MATERIAL


def _file_type(ext: str) -> str:
    return {"pdf": "pdf", "txt": "txt", "md": "md", "markdown": "md", "pptx": "pptx", "doc": "doc", "docx": "docx"}.get(ext.lower().lstrip("."), "other")


def _md5_of(path: Path) -> str:
    h = hashlib.md5()
    with path.open("rb") as f:
        for block in iter(lambda: f.read(1 << 16), b""):
            h.update(block)
    return h.hexdigest()


def _read_text(full_path: Path, ext: str) -> str:
    if ext == ".pdf":
        try:
            from pypdf import PdfReader
            reader = PdfReader(str(full_path))
            return "\n".join((p.extract_text() or "") for p in reader.pages)
        except Exception as e:
            raise ValueError(f"PDF 解析失败：{e}")
    if ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(str(full_path))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                        parts.append(shape.text_frame.text.strip())
                    if getattr(shape, "has_table", False):
                        for row in shape.table.rows:
                            cells = [c.text.strip() for c in row.cells]
                            if any(cells):
                                parts.append(" | ".join(cells))
            text = "\n".join(parts).strip()
            if not text:
                raise ValueError("PPTX 中未提取到文字内容（可能全部是图片/图形）")
            return text
        except ValueError:
            raise
        except Exception as e:
            raise ValueError(f"PPTX 解析失败：{e}")
    if ext == ".docx":
        try:
            from docx import Document as DocxDocument
            from docx.oxml.ns import qn

            d = DocxDocument(str(full_path))
            parts = [p.text.strip() for p in d.paragraphs if p.text.strip()]
            for table in d.tables:
                for row in table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))
            # 正文段落/表格不覆盖文本框、形状内的文字，这里全量扫描 w:t 兜底
            if not parts:
                parts = [
                    t.text.strip()
                    for t in d.element.body.iter(qn("w:t"))
                    if t.text and t.text.strip()
                ]
            text = "\n".join(parts).strip()
            if not text:
                raise ValueError("Word 文档中未提取到文字内容（可能全部是图片/扫描件）")
            return text
        except ValueError:
            raise
        except Exception as e:
            # 扩展名是 .docx 但内容可能是旧版 .doc（二进制），尝试兼容解析
            try:
                return _read_legacy_doc(full_path)
            except Exception:
                pass
            raise ValueError(f"Word 文档解析失败：{e}")
    if ext == ".doc":
        return _read_legacy_doc(full_path)
    try:
        return full_path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return full_path.read_text(encoding="gbk", errors="ignore")


def _strip_rtf(data: bytes) -> str:
    """粗略解析 RTF：还原 \\uN? 转义、去掉控制字与花括号（供 .doc 兼容用）。"""
    import re

    def _u(m):
        try:
            n = int(m.group(1))
            return chr(n) if -65535 <= n <= 65535 else ""
        except Exception:
            return ""

    text = data.decode("latin-1", errors="ignore")
    text = re.sub(r"\\u(-?\d+)\?", _u, text)
    text = re.sub(r"\\[a-zA-Z]+-?\d* ?", "", text)
    text = text.replace("{", "").replace("}", "")
    return text.strip()


def _decode_ole_word(data: bytes) -> str:
    """从 .doc 的 WordDocument 流粗略提取可见文本（启发式，供 RAG 检索用）。"""
    import re

    s = data.decode("utf-16-le", errors="ignore")
    out = []
    for ch in s:
        out.append(ch if ch in "\r\n\t" or ch.isprintable() else " ")
    text = "".join(out)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _read_legacy_doc(full_path: Path) -> str:
    """兼容旧版 .doc：先判断是否 RTF，再尝试 OLE 二进制提取，均失败给出引导。"""
    raw = full_path.read_bytes()
    if raw.lstrip().startswith(b"{\\rtf"):
        text = _strip_rtf(raw)
        if text:
            return text
    try:
        import olefile

        if olefile.isOleFile(str(full_path)):
            ole = olefile.OleFileIO(str(full_path))
            try:
                data = ole.openstream("WordDocument").read()
            finally:
                ole.close()
            text = _decode_ole_word(data)
            if text:
                return text
    except Exception:
        pass
    raise ValueError("暂不支持解析旧版 .doc（Word 97-2003 二进制）文件，请用 Word 另存为 .docx 后重新上传")


def _chunk_text(text: str) -> List[str]:
    if not text.strip():
        return []
    try:
        from langchain_text_splitters import RecursiveCharacterTextSplitter
    except Exception:
        from langchain.text_splitter import RecursiveCharacterTextSplitter  # type: ignore
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=settings.RAG_CHUNK_SIZE,
        chunk_overlap=settings.RAG_CHUNK_OVERLAP,
        separators=["\n\n", "\n", "。", "！", "？", ".", ",", " ", ""],
    )
    return splitter.split_text(text)


def save_upload(
    owner_id: int,
    file: UploadFile,
    db: Session,
    *,
    category: str | None = None,
    description: str | None = None,
) -> DocumentUploaded:
    raw = file.filename or "unnamed"
    ext = os.path.splitext(raw)[1].lower()
    if ext not in _ALLOWED_EXT:
        raise ValueError(f"不支持的文件类型：{ext or '无'}，支持 {sorted(_ALLOWED_EXT)}")

    storage = _ensure_dirs()
    safe_name = f"{datetime.now():%Y%m%d}_{uuid.uuid4().hex[:8]}_{os.path.basename(raw)}"
    target = storage / f"u{owner_id}" / safe_name
    target.parent.mkdir(parents=True, exist_ok=True)
    with target.open("wb") as fout:
        shutil.copyfileobj(file.file, fout)
    file_size = target.stat().st_size
    max_bytes = settings.DOC_MAX_MB * 1024 * 1024
    if file_size > max_bytes:
        target.unlink(missing_ok=True)
        raise ValueError(f"文件过大：{file_size // 1024 // 1024}MB，上限 {settings.DOC_MAX_MB}MB")

    # 内容去重：同用户相同 MD5 视为重复文件
    digest = _md5_of(target)
    existing = db.scalar(
        select(Document).where(Document.owner_id == owner_id, Document.md5 == digest)
    )
    if existing is not None:
        target.unlink(missing_ok=True)
        raise ValueError(f"相同内容的文件已存在：{existing.title}")

    try:
        text = _read_text(target, ext)
    except Exception:
        target.unlink(missing_ok=True)
        raise

    chunks = _chunk_text(text)

    doc = Document(
        owner_id=owner_id,
        title=raw,
        file_type=_file_type(ext),
        file_path=str(target),
        file_size=file_size,
        md5=digest,
        category=_category_of(raw, category),
        chroma_collection=user_collection_name(owner_id),
        chunk_count=len(chunks),
        description=description,
    )
    db.add(doc)
    db.flush()

    if chunks:
        ids = [f"d{doc.id}_c{i}" for i in range(len(chunks))]
        metas = [
            {"doc_id": doc.id, "doc_title": doc.title, "category": doc.category, "chunk_index": i}
            for i in range(len(chunks))
        ]
        vs = get_user_vectorstore(owner_id)
        vs.add_texts(texts=chunks, metadatas=metas, ids=ids)

    db.commit()
    db.refresh(doc)
    return DocumentUploaded(id=doc.id, title=doc.title, chunk_count=doc.chunk_count, category=doc.category)


def list_my_documents(
    owner_id: int,
    db: Session,
    *,
    category: str | None = None,
    keyword: str | None = None,
    page: int = 1,
    page_size: int = 20,
) -> tuple[List[DocumentOut], int]:
    stmt = select(Document).where(Document.owner_id == owner_id)
    if category:
        stmt = stmt.where(Document.category == category)
    if keyword:
        stmt = stmt.where(
            or_(
                Document.title.contains(keyword),
                Document.description.contains(keyword),
            )
        )
    total = db.scalar(select(func.count()).select_from(stmt.subquery())) or 0
    rows = db.scalars(
        stmt.order_by(Document.created_at.desc())
        .offset((page - 1) * page_size)
        .limit(page_size)
    ).all()
    return [DocumentOut.model_validate(r) for r in rows], int(total)


def get_document(owner_id: int, doc_id: int, db: Session) -> Document | None:
    return db.scalar(select(Document).where(Document.id == doc_id, Document.owner_id == owner_id))


def delete_document(owner_id: int, doc_id: int, db: Session) -> bool:
    doc = get_document(owner_id, doc_id, db)
    if doc is None:
        return False
    try:
        vs = get_user_vectorstore(owner_id, create_if_missing=False)
        col = vs._collection  # noqa: SLF001
        try:
            col.delete(where={"doc_id": doc.id})
        except Exception:
            res = col.get(where={}, include=[])
            prefix = f"d{doc.id}_c"
            to_del = [x for x in (res.get("ids") or []) if str(x).startswith(prefix)]
            if to_del:
                col.delete(ids=to_del)
    except Exception:
        pass
    try:
        Path(doc.file_path).unlink(missing_ok=True)
    except Exception:
        pass
    db.delete(doc)
    db.commit()
    return True


def _update_chroma_title(owner_id: int, doc: Document, new_title: str) -> None:
    """重命名后同步向量库中的 doc_title 元数据（失败不影响主流程）。"""
    try:
        vs = get_user_vectorstore(owner_id, create_if_missing=False)
        col = vs._collection  # noqa: SLF001
        res = col.get(where={"doc_id": doc.id}, include=["metadatas"])
        ids = res.get("ids") or []
        metas = res.get("metadatas") or []
        if not ids:
            return
        col.update(
            ids=ids,
            metadatas=[{**(m or {}), "doc_title": new_title} for m in metas],
        )
    except Exception:
        pass


def update_document(
    owner_id: int,
    doc_id: int,
    db: Session,
    *,
    title: str | None = None,
    description: str | None = None,
) -> DocumentOut | None:
    doc = get_document(owner_id, doc_id, db)
    if doc is None:
        return None
    if title is not None:
        doc.title = title
        _update_chroma_title(owner_id, doc, title)
    if description is not None:
        doc.description = description
    db.add(doc)
    db.commit()
    db.refresh(doc)
    return DocumentOut.model_validate(doc)


def preview_text(path: str) -> str:
    """读取文档文本用于前端预览（PDF 走文件流，不在此处理）。"""
    ext = os.path.splitext(path)[1]
    return _read_text(Path(path), ext)
